import os
from pptx import Presentation
from pptx.util import Inches
import pandas as pd

# Ensure the output directory exists
os.makedirs("outputs", exist_ok=True)

df = pd.read_excel("outputs/processed_stats.xlsx")
top_attendance = df.sort_values(by="Total Attendance", ascending=False).head(5)

prs = Presentation()

# Slide 1: Title
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
slide_1.shapes.title.text = "MMA Club Member Stats"
slide_1.placeholders[1].text = "Summary of attendance, sparring performance, and efficiency"

# Slide 2: Table with top performers
slide_2 = prs.slides.add_slide(prs.slide_layouts[6])
top_table_data = top_attendance[["Name", "Total Attendance", "Efficiency Score"]]
rows, cols = top_table_data.shape[0] + 1, top_table_data.shape[1]
table_shape = slide_2.shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), Inches(9), Inches(3))
table = table_shape.table
for i, col in enumerate(top_table_data.columns):
    table.cell(0, i).text = str(col)
for i, row in enumerate(top_table_data.itertuples(index=False)):
    for j, value in enumerate(row):
        table.cell(i + 1, j).text = str(value)

# Slide 3: Performance commentary
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
slide_3.shapes.title.text = "Performance Analysis"
slide_3.placeholders[1].text = (
    "Efficiency Score is calculated as: (Wins - Losses) × Endurance\n"
    "This metric combines sparring results with physical stamina.\n\n"
    f"Top performer: {df.loc[df['Efficiency Score'].idxmax(), 'Name']}\n"
    f"Highest Efficiency Score: {df['Efficiency Score'].max():.2f}"
)

# Save presentation file
prs.save("outputs/MMA_Club_Stats_Presentation.pptx")
print("Saved successfully presentation file in 'outputs/MMA_Club_Stats_Presentation.pptx'")